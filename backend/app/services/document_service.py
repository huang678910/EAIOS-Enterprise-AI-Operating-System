"""文档服务 — 上传 / 解析 / Chunk / Embedding / 存储"""

import uuid
import logging
from pathlib import Path
from typing import BinaryIO

from fastapi import UploadFile
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select

from app.config import get_settings
from app.models.document import Document, DocumentStatus
from app.models.document_chunk import DocumentChunk
from app.core.exceptions import NotFoundException, BadRequestException

logger = logging.getLogger(__name__)
settings = get_settings()

# ---- 文件解析器 ----

def parse_pdf(filepath: str) -> str:
    """PDF 解析 — 优先使用 PyMuPDF (fitz)，回退到 pypdf"""
    # Try PyMuPDF first (better text extraction + image support)
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        parts = []
        for page in doc:
            text = page.get_text()
            if text.strip():
                parts.append(text)
            # Extract embedded images
            for img_index, img in enumerate(page.get_images(full=True)):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    img_ext = base_image["ext"]
                    if img_bytes and len(img_bytes) > 1024:  # Skip tiny images
                        import base64
                        b64 = base64.b64encode(img_bytes).decode()
                        parts.append(f"[Image: page_{page.number + 1}_img_{img_index}.{img_ext}]\n(data:image/{img_ext};base64,{b64[:100]}...)")
                except Exception:
                    pass
        doc.close()
        return "\n\n".join(parts) if parts else "(No extractable text)"
    except ImportError:
        pass

    # Fallback to pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except ImportError:
        raise ImportError("pypdf or PyMuPDF not installed")

def parse_docx(filepath: str) -> str:
    try:
        from docx import Document as DocxDoc
        doc = DocxDoc(filepath)
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        raise ImportError("python-docx not installed")

def parse_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        from pptx.shapes.picture import Picture
        prs = Presentation(filepath)
        text = []
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text_frame.text)
                # Check for images
                if hasattr(shape, "image"):
                    try:
                        img = shape.image
                        text.append(f"[Image: slide_{slide_num + 1}_{shape.name}.{img.content_type.split('/')[-1]}]")
                    except Exception:
                        pass
                # Check for tables
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        text.append("| " + " |\n| ".join(r for r in rows) + " |")
        return "\n".join(text)
    except ImportError:
        raise ImportError("python-pptx not installed")

def parse_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def parse_csv(filepath: str) -> str:
    """解析 CSV 文件为结构化文本"""
    try:
        import pandas as pd
        df = pd.read_csv(filepath)
        # 限制行数
        if len(df) > 1000:
            df = df.head(1000)
        return df.to_markdown(index=False)
    except ImportError:
        raise ImportError("pandas not installed for CSV parsing")
    except Exception as e:
        raise ValueError(f"Failed to parse CSV: {e}")


def parse_xlsx(filepath: str) -> str:
    """解析 Excel (.xlsx/.xls) 文件为结构化文本"""
    try:
        import pandas as pd
        # Read all sheets
        xl = pd.ExcelFile(filepath)
        parts = []
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(filepath, sheet_name=sheet_name)
            if df.empty:
                continue
            if len(df) > 500:
                df = df.head(500)
            parts.append(f"## Sheet: {sheet_name}\n{df.to_markdown(index=False)}")
        return "\n\n".join(parts) if parts else "(Empty spreadsheet)"
    except ImportError:
        raise ImportError("pandas + openpyxl required for Excel parsing. Install: pip install pandas openpyxl")
    except Exception as e:
        raise ValueError(f"Failed to parse Excel: {e}")


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,
    ".txt": parse_txt,
    ".md": parse_txt,
    ".markdown": parse_txt,
    ".csv": parse_csv,
    ".xlsx": parse_xlsx,
    ".xls": parse_xlsx,
}

SUPPORTED_TYPES = set(PARSERS.keys())


# ---- 文档服务 ----

class DocumentService:
    def __init__(self, db: AsyncSession):
        self.db = db

    async def create_document(
        self, workspace_id: str, filename: str, file_type: str,
        file_size: int = 0, source_type: str = "upload",
    ) -> Document:
        """Create a document record (for connectors / transcription)"""
        doc = Document(
            workspace_id=workspace_id,
            filename=filename,
            file_type=file_type,
            file_size=file_size,
            status=DocumentStatus.READY,
            source_type=source_type,
        )
        self.db.add(doc)
        await self.db.flush()
        await self.db.refresh(doc)
        return doc

    async def list_by_workspace(self, workspace_id: str) -> list[Document]:
        result = await self.db.execute(
            select(Document)
            .where(Document.workspace_id == workspace_id)
            .order_by(Document.created_at.desc())
        )
        return list(result.scalars().all())

    async def upload(self, workspace_id: str, file: UploadFile) -> dict:
        # 验证文件类型
        filename = file.filename or "unknown"
        ext = Path(filename).suffix.lower()
        if ext not in SUPPORTED_TYPES:
            raise BadRequestException(
                f"Unsupported file type '{ext}'. Supported: {', '.join(sorted(SUPPORTED_TYPES))}"
            )

        # 流式写入文件到磁盘（避免大文件撑爆内存）
        upload_dir = Path(settings.UPLOAD_DIR)
        upload_dir.mkdir(parents=True, exist_ok=True)
        file_id = uuid.uuid4()
        filepath = upload_dir / f"{file_id}_{filename}"

        CHUNK_SIZE = 1024 * 1024  # 1MB per chunk
        total_size = 0
        with open(filepath, "wb") as f:
            while chunk := await file.read(CHUNK_SIZE):
                f.write(chunk)
                total_size += len(chunk)

        # 创建文档记录
        doc = Document(
            id=file_id,
            workspace_id=workspace_id,
            filename=filename,
            file_type=ext.lstrip("."),
            file_size=total_size,
            status=DocumentStatus.PROCESSING,
            file_path=str(filepath),
        )
        self.db.add(doc)
        await self.db.flush()

        try:
            # 解析 → Chunk → Embedding → 存储
            await self._process_document(doc)
        except Exception as e:
            logger.error(f"Document processing failed: {e}")
            doc.status = DocumentStatus.ERROR
            doc.error_message = str(e)
            await self.db.flush()
            return {
                "id": str(doc.id),
                "filename": doc.filename,
                "status": "error",
                "message": str(e),
            }

        return {
            "id": str(doc.id),
            "filename": doc.filename,
            "status": "ready",
            "message": f"Document processed successfully. {doc.chunk_count} chunks created.",
        }

    async def delete(self, document_id: str) -> None:
        result = await self.db.execute(
            select(Document).where(Document.id == document_id)
        )
        doc = result.scalar_one_or_none()
        if not doc:
            raise NotFoundException("Document not found")

        # 删除磁盘文件
        if doc.file_path:
            try:
                Path(doc.file_path).unlink(missing_ok=True)
            except Exception:
                pass

        await self.db.delete(doc)

    async def _process_document(self, doc: Document) -> None:
        """核心处理流水线：解析 → 分段 → 向量化 → 存储"""
        from app.services.embedding_service import embed_texts

        # 1. 解析文件
        filepath = doc.file_path
        if not filepath:
            raise ValueError("Document has no file path")

        ext = Path(filepath).suffix.lower()
        parser = PARSERS.get(ext)
        if not parser:
            raise ValueError(f"No parser for extension {ext}")

        text = parser(filepath)
        if not text.strip():
            raise ValueError("No text could be extracted from the document")

        # 2. Chunk 分割
        chunks = self._split_text(text)

        # 3. 生成 Embedding
        chunk_texts = [c["content"] for c in chunks]
        embeddings = await embed_texts(chunk_texts)

        # 4. 批量存入 document_chunks
        for i, chunk in enumerate(chunks):
            chunk_record = DocumentChunk(
                document_id=doc.id,
                chunk_index=i,
                content=chunk["content"],
                metadata_=chunk.get("metadata", {}),
                embedding=embeddings[i],
            )
            self.db.add(chunk_record)

        doc.status = DocumentStatus.READY
        doc.chunk_count = len(chunks)
        await self.db.flush()

    def _split_text(self, text: str, chunk_size: int = 512, overlap: int = 64) -> list[dict]:
        """使用 LangChain RecursiveCharacterTextSplitter 分割文本"""
        from langchain.text_splitter import RecursiveCharacterTextSplitter

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", "。", ".", " ", ""],
            length_function=len,
        )
        raw_chunks = splitter.split_text(text)

        chunks = []
        for i, chunk_text in enumerate(raw_chunks):
            if chunk_text.strip():
                chunks.append({
                    "content": chunk_text.strip(),
                    "metadata": {
                        "chunk_index": i,
                        "char_count": len(chunk_text),
                    },
                })
        return chunks
