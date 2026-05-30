"""
MCP Server wrapping LightRAG for enterprise knowledge base.
Supports: PDF, DOCX, PPTX, TXT, Markdown ingest + hybrid vector/graph query.
"""
import os
import sys
import json
import asyncio
import logging
from pathlib import Path
from typing import Any

from mcp.server import Server, InitializationOptions, NotificationOptions
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent, ServerCapabilities, ToolsCapability

# --- LightRAG ---
from lightrag import LightRAG, QueryParam
from lightrag.llm.openai import openai_complete_if_cache, openai_embed
from lightrag.utils import EmbeddingFunc

# --- Config ---
WORKING_DIR = os.environ.get("KNOWLEDGE_WORKING_DIR", os.path.expanduser("~/.claude/knowledge_base"))
LLM_API_KEY = os.environ.get("DEEPSEEK_API_KEY", os.environ.get("OPENAI_API_KEY", ""))
LLM_BASE_URL = os.environ.get("OPENAI_BASE_URL", "https://api.deepseek.com/v1")
LLM_MODEL = os.environ.get("LLM_MODEL", "deepseek-chat")
# Embedding: "local" = free local model, or set a model name for API-based embeddings
EMBEDDING_MODE = os.environ.get("EMBEDDING_MODE", "local")  # "local" or "api"
EMBEDDING_MODEL = os.environ.get("EMBEDDING_MODEL", "paraphrase-multilingual-MiniLM-L12-v2")
EMBEDDING_API_KEY = os.environ.get("EMBEDDING_API_KEY", LLM_API_KEY)
EMBEDDING_BASE_URL = os.environ.get("EMBEDDING_BASE_URL", LLM_BASE_URL)

logging.basicConfig(level=logging.INFO, stream=sys.stderr)
logger = logging.getLogger("knowledge-mcp")


# --- LLM & Embedding Functions ---

async def llm_model_func(prompt, system_prompt=None, history_messages=None, **kwargs):
    return await openai_complete_if_cache(
        LLM_MODEL, prompt, system_prompt=system_prompt,
        history_messages=history_messages, api_key=LLM_API_KEY,
        base_url=LLM_BASE_URL, **kwargs
    )


# Local embedding model (lazy loaded)
_local_embed_model = None

def _get_local_embed_model():
    global _local_embed_model
    if _local_embed_model is None:
        from sentence_transformers import SentenceTransformer
        logger.info(f"Loading local embedding model: {EMBEDDING_MODEL}")
        _local_embed_model = SentenceTransformer(EMBEDDING_MODEL)
    return _local_embed_model


async def embedding_func(texts: list[str]) -> list[list[float]]:
    if EMBEDDING_MODE == "local":
        model = _get_local_embed_model()
        embeddings = model.encode(texts, normalize_embeddings=True)
        return embeddings.tolist()
    else:
        return await openai_embed(
            texts, model=EMBEDDING_MODEL,
            api_key=EMBEDDING_API_KEY, base_url=EMBEDDING_BASE_URL
        )


# --- Embedding dimension: 384 for MiniLM, 1536 for OpenAI ---
EMBEDDING_DIM = 384 if EMBEDDING_MODE == "local" else 1536

# Initialize LightRAG
rag = LightRAG(
    working_dir=WORKING_DIR,
    llm_model_func=llm_model_func,
    embedding_func=EmbeddingFunc(
        embedding_dim=EMBEDDING_DIM,
        max_token_size=8192,
        func=embedding_func,
    ),
)

# --- Document Parsers ---

def parse_pdf(filepath: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except ImportError:
        raise ImportError("pypdf not installed. Run: pip install pypdf")

def parse_docx(filepath: str) -> str:
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")

def parse_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text_frame.text)
        return "\n".join(text)
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

def parse_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def parse_md(filepath: str) -> str:
    return parse_txt(filepath)

PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,
    ".txt": parse_txt,
    ".md": parse_md,
    ".markdown": parse_md,
}

# --- MCP Server ---

server = Server("knowledge-mcp")


@server.list_tools()
async def list_tools() -> list[Tool]:
    return [
        Tool(
            name="kb_ingest_document",
            description="Ingest a document into the knowledge base. Supports PDF, DOCX, PPTX, TXT, Markdown files.",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Absolute path to the document file to ingest."
                    },
                    "title": {
                        "type": "string",
                        "description": "Optional title for the document. If not provided, the filename is used."
                    }
                },
                "required": ["filepath"]
            }
        ),
        Tool(
            name="kb_search",
            description="Search the knowledge base using hybrid vector + knowledge graph retrieval. Returns relevant information from all ingested documents.",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "The search query or question."
                    },
                    "mode": {
                        "type": "string",
                        "enum": ["hybrid", "local", "global", "naive"],
                        "description": "Search mode: hybrid (vector+graph, recommended), local (entity-focused), global (relationship-focused), naive (keyword+vector)."
                    },
                    "top_k": {
                        "type": "integer",
                        "description": "Number of results to return. Default 10."
                    }
                },
                "required": ["query"]
            }
        ),
        Tool(
            name="kb_list_documents",
            description="List all documents currently in the knowledge base.",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        ),
        Tool(
            name="kb_delete_document",
            description="Remove a document from the knowledge base by its ID.",
            inputSchema={
                "type": "object",
                "properties": {
                    "doc_id": {
                        "type": "string",
                        "description": "The document ID to delete (shown in kb_list_documents)."
                    }
                },
                "required": ["doc_id"]
            }
        ),
        Tool(
            name="kb_stats",
            description="Get statistics about the knowledge base (document count, entity count, relation count).",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict[str, Any]) -> list[TextContent]:
    try:
        if name == "kb_ingest_document":
            filepath = arguments["filepath"]
            title = arguments.get("title", Path(filepath).stem)

            if not os.path.isfile(filepath):
                return [TextContent(type="text", text=f"Error: File not found: {filepath}")]

            ext = Path(filepath).suffix.lower()
            if ext not in PARSERS:
                supported = ", ".join(PARSERS.keys())
                return [TextContent(type="text", text=f"Error: Unsupported file type '{ext}'. Supported: {supported}")]

            text = PARSERS[ext](filepath)
            if not text.strip():
                return [TextContent(type="text", text=f"Error: No text extracted from {filepath}")]

            await rag.ainsert(text, file_paths=str(filepath))
            doc_count = len(rag.get_docs_by_status("processed"))
            return [TextContent(
                type="text",
                text=f"Document '{title}' ingested successfully.\n"
                     f"  File: {filepath}\n"
                     f"  Characters: {len(text):,}\n"
                     f"  Knowledge base now has {doc_count} document(s).\n"
                     f"  Use kb_search to query the knowledge base."
            )]

        elif name == "kb_search":
            query = arguments["query"]
            mode = arguments.get("mode", "hybrid")
            top_k = arguments.get("top_k", 10)

            param = QueryParam(mode=mode, top_k=top_k)
            result = await rag.aquery(query, param=param)
            return [TextContent(type="text", text=result)]

        elif name == "kb_list_documents":
            docs = rag.get_docs_by_status("processed")
            if not docs:
                return [TextContent(type="text", text="Knowledge base is empty. Use kb_ingest_document to add documents.")]

            lines = [f"Knowledge Base Documents ({len(docs)} total):", ""]
            for doc_id, info in docs.items():
                file_paths = info.get("file_paths", [])
                created = info.get("created_at", "unknown")
                lines.append(f"  - ID: {doc_id}")
                lines.append(f"    Files: {', '.join(file_paths) if file_paths else 'N/A'}")
                lines.append(f"    Created: {created}")
                lines.append("")
            return [TextContent(type="text", text="\n".join(lines))]

        elif name == "kb_delete_document":
            doc_id = arguments["doc_id"]
            await rag.adelete_by_doc_id(doc_id)
            return [TextContent(type="text", text=f"Document '{doc_id}' deleted.")]

        elif name == "kb_stats":
            docs = rag.get_docs_by_status("processed")
            # Count entities and relations from the knowledge graph
            try:
                kg = rag.get_knowledge_graph()
                entity_count = len(kg.get("nodes", {}))
                relation_count = len(kg.get("edges", {}))
            except Exception:
                entity_count = "N/A"
                relation_count = "N/A"

            stats = (
                f"Knowledge Base Statistics:\n"
                f"  Documents: {len(docs)}\n"
                f"  Entities (nodes): {entity_count}\n"
                f"  Relations (edges): {relation_count}\n"
                f"  Storage directory: {os.path.abspath(WORKING_DIR)}"
            )
            return [TextContent(type="text", text=stats)]

        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]

    except Exception as e:
        logger.error(f"Tool '{name}' error: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error: {str(e)}")]


async def main():
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="knowledge-mcp",
                server_version="1.0.0",
                capabilities=ServerCapabilities(
                    tools=ToolsCapability(),
                ),
            ),
        )


if __name__ == "__main__":
    asyncio.run(main())
